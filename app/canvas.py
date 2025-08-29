from __future__ import annotations
import io, csv
from typing import Any, Dict, List, Optional, Tuple
import httpx
from bs4 import BeautifulSoup
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

class CanvasClient:
    def __init__(self, base_url: str, token: str):
        self.base = (base_url or "https://canvas.instructure.com/").rstrip("/")
        self.token = token

    def _h(self) -> Dict[str, str]:
        return {"Authorization": f"Bearer {self.token}"}

    def ping(self) -> None:
        url = f"{self.base}/api/v1/courses?per_page=1"
        with httpx.Client(timeout=30, follow_redirects=True) as s:
            r = s.get(url, headers=self._h())
            if r.status_code == 401:
                raise RuntimeError(r.text)
            r.raise_for_status()

    def list_courses(self) -> List[Dict[str, Any]]:
        url = f"{self.base}/api/v1/courses?per_page=100&enrollment_state=active"
        with httpx.Client(timeout=30, follow_redirects=True) as s:
            r = s.get(url, headers=self._h()); r.raise_for_status()
            out = []
            for c in r.json():
                out.append({"id": c.get("id"), "name": c.get("name") or c.get("course_code") or f"Course {c.get('id')}"})
            return out

    def list_modules(self, course_id: int) -> List[Dict[str, Any]]:
        url = f"{self.base}/api/v1/courses/{course_id}/modules?per_page=100"
        with httpx.Client(timeout=30, follow_redirects=True) as s:
            r = s.get(url, headers=self._h()); r.raise_for_status()
            mods = r.json()
        out: List[Dict[str, Any]] = []
        for m in mods:
            mid = m.get("id"); name = m.get("name") or f"Module {mid}"
            items_url = f"{self.base}/api/v1/courses/{course_id}/modules/{mid}/items?per_page=200"
            with httpx.Client(timeout=30, follow_redirects=True) as s:
                ri = s.get(items_url, headers=self._h()); ri.raise_for_status()
                raws = ri.json()
            items: List[Dict[str, Any]] = []
            for it in raws:
                typ = it.get("type")
                if typ == "Page":
                    items.append({"id": it.get("id"), "type": "Page", "title": it.get("title"), "page_url": it.get("page_url")})
                elif typ == "File":
                    fid = it.get("content_id") or it.get("id")
                    items.append({"id": it.get("id"), "type": "File", "title": it.get("title"), "file_id": fid})
                elif typ == "Assignment":
                    aid = it.get("content_id") or it.get("id")
                    items.append({"id": it.get("id"), "type": "Assignment", "title": it.get("title"), "assignment_id": aid})
            out.append({"id": mid, "name": name, "items": items})
        return out

    def get_page_body(self, course_id: int, page_url: str) -> str:
        url = f"{self.base}/api/v1/courses/{course_id}/pages/{page_url}"
        with httpx.Client(timeout=30, follow_redirects=True) as s:
            r = s.get(url, headers=self._h())
            if r.status_code == 404: return ""
            r.raise_for_status()
            body = (r.json().get("body") or "")
        soup = BeautifulSoup(body, "html.parser")
        return (soup.get_text(separator="\n") or "").strip()

    def get_assignment_text(self, course_id: int, assignment_id: int) -> str:
        url = f"{self.base}/api/v1/courses/{course_id}/assignments/{assignment_id}"
        with httpx.Client(timeout=30, follow_redirects=True) as s:
            r = s.get(url, headers=self._h())
            if r.status_code == 404: return ""
            r.raise_for_status()
            desc = r.json().get("description") or ""
        soup = BeautifulSoup(desc, "html.parser")
        return (soup.get_text(separator="\n") or "").strip()

    def get_file_text(self, file_id: int) -> Tuple[str, str]:
        """Return (text, warning). Tries PDF/DOCX/PPTX/CSV/TSV; otherwise UTF-8 text."""
        meta_url = f"{self.base}/api/v1/files/{file_id}"
        with httpx.Client(timeout=60, follow_redirects=True) as s:
            r = s.get(meta_url, headers=self._h())
            if r.status_code == 404:
                return "", f"File {file_id}: not found"
            r.raise_for_status()
            meta = r.json()
            dl = meta.get("url") or meta.get("download_url") or meta.get("preview_url")
            fname = meta.get("display_name") or meta.get("filename") or f"file_{file_id}"
            if not dl:
                return "", f"File {file_id} ({fname}): no download_url"
            # Try without headers first (presigned). If it fails, try with headers.
            rf = s.get(dl)
            if rf.status_code >= 400:
                rf = s.get(dl, headers=self._h())
            if rf.status_code >= 400:
                return "", f"{fname}: download failed ({rf.status_code})"
            content = rf.content

        low = (fname or "").lower()
        try:
            if low.endswith(".pdf"):
                pdf = PdfReader(io.BytesIO(content))
                return "\n".join([(p.extract_text() or "") for p in pdf.pages]).strip(), ""
            if low.endswith(".docx"):
                doc = Document(io.BytesIO(content))
                return "\n".join([p.text for p in doc.paragraphs]).strip(), ""
            if low.endswith(".pptx"):
                prs = Presentation(io.BytesIO(content))
                slides = []
                for sld in prs.slides:
                    segs = []
                    for shp in sld.shapes:
                        if hasattr(shp, "text"):
                            segs.append(shp.text)
                    slides.append("\n".join(segs))
                return "\n\n".join(slides).strip(), ""
            if low.endswith(".csv") or low.endswith(".tsv"):
                import io as _io
                text = content.decode("utf-8", errors="ignore")
                delim = "," if low.endswith(".csv") else "\t"
                rows = [" | ".join(row) for row in csv.reader(_io.StringIO(text), delimiter=delim)]
                return "\n".join(rows), ""
            # Fallback: plain text
            return content.decode("utf-8", errors="ignore").strip(), ""
        except Exception as e:
            return "", f"{fname}: extract failed ({e})"

    def _b(self, v: bool) -> str: return "true" if v else "false"

    # --- Classic Quiz creation ---
    def create_quiz(self, course_id: int, title: str, settings: dict | None = None) -> dict:
        from urllib.parse import urlencode
        fields = [("quiz[title]", str(title or "Generated Quiz"))]
        s = settings or {}
        if s.get("description"):      fields.append(("quiz[description]", str(s["description"])))
        if "published" in s:          fields.append(("quiz[published]", self._b(bool(s["published"]))))
        if "time_limit" in s and s["time_limit"] is not None:
                                      fields.append(("quiz[time_limit]", str(int(s["time_limit"]))))
        if s.get("due_at"):           fields.append(("quiz[due_at]", str(s["due_at"])))
        if "shuffle_answers" in s:    fields.append(("quiz[shuffle_answers]", self._b(bool(s["shuffle_answers"]))))
        if s.get("scoring_policy"):   fields.append(("quiz[scoring_policy]", str(s["scoring_policy"])))
        body = urlencode(fields, doseq=True).encode("utf-8")
        headers = {"Authorization": f"Bearer {self.token}", "Content-Type": "application/x-www-form-urlencoded"}
        url = f"{self.base}/api/v1/courses/{course_id}/quizzes"
        with httpx.Client(timeout=30, follow_redirects=True) as scli:
            r = scli.post(url, content=body, headers=headers)
        if r.status_code >= 400:
            try: msg = r.json()
            except Exception: msg = r.text
            raise RuntimeError(f"Canvas error creating quiz: {msg}")
        return r.json()

    def create_quiz_question(self, course_id: int, quiz_id: int, q: dict) -> dict:
        from urllib.parse import urlencode
        tmap = {"mcq":"multiple_choice_question","truefalse":"true_false_question","short":"short_answer_question","fillblank":"fill_in_the_blank_question"}
        qtype = tmap.get((q.get("type") or "short").lower(), "short_answer_question")
        prompt = str(q.get("prompt") or "Question")
        name = (prompt[:50] + ("…" if len(prompt) > 50 else ""))
        points = str(int(q.get("points") or 1))
        fields: List[tuple[str,str]] = [
            ("question[question_name]", name),
            ("question[question_text]", prompt),
            ("question[points_possible]", points),
            ("question[question_type]", qtype),
        ]
        if qtype == "multiple_choice_question":
            choices = list(q.get("choices") or [])
            correct_idx = q.get("answer")
            for i, choice in enumerate(choices):
                fields.append(("question[answers][][text]", str(choice)))
                fields.append(("question[answers][][weight]", "100" if correct_idx == i else "0"))
        elif qtype == "true_false_question":
            ans = bool(q.get("answer"))
            fields.append(("question[answers][][text]", "True"))
            fields.append(("question[answers][][weight]", "100" if ans else "0"))
            fields.append(("question[answers][][text]", "False"))
            fields.append(("question[answers][][weight]", "0" if ans else "100"))
        elif qtype == "fill_in_the_blank_question":
            ans = str(q.get("answer") or "").strip()
            if ans:
                fields.append(("question[answers][][text]", ans))
                fields.append(("question[answers][][weight]", "100"))
        body = urlencode(fields, doseq=True).encode("utf-8")
        headers = {"Authorization": f"Bearer {self.token}", "Content-Type": "application/x-www-form-urlencoded"}
        url = f"{self.base}/api/v1/courses/{course_id}/quizzes/{quiz_id}/questions"
        with httpx.Client(timeout=30, follow_redirects=True) as scli:
            r = scli.post(url, content=body, headers=headers)
        if r.status_code >= 400:
            try: msg = r.json()
            except Exception: msg = r.text
            raise RuntimeError(f"Canvas error creating question: {msg}")
        return r.json()
